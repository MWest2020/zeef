#!/usr/bin/env python3
# SPDX-License-Identifier: EUPL-1.2
# role: tool
#
# presentation/pptx-build/build_pptx.py — bouwt verkenning-bzk.pptx uit de reveal.js-deck.
#
# Wat: recreëert de 9 slides van presentation/verkenning-bzk.html als een native,
# bewerkbare PowerPoint. Tekst, tabellen, kaarten en de trechter zijn echte PPTX-objecten
# (bewerkbaar in PowerPoint). De twee SVG-diagrammen (double-diamond, venn) en de twee
# echte rapport-embeds zijn vooraf gerasterd naar PNG (zie img/) en als afbeelding ingesloten.
# Speaker notes uit de <aside class="notes"> zijn overgenomen in de notitievelden.
#
# Writes: ../verkenning-bzk.pptx
# Idempotent: ja (overschrijft de output volledig).
# Requires: python-pptx, Pillow (PIL — picture_fit leest beeldafmetingen); vooraf gegenereerde
#   PNG's in img/ (logo, diamond, venn, converge, discover).
#
# Usage:
#   python build_pptx.py

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
OUT = os.path.join(HERE, "..", "verkenning-bzk.pptx")

# ---- Palet (1-op-1 uit het CSS-thema) ----
INK        = RGBColor(0x1a, 0x1a, 0x1a)
G700       = RGBColor(0x3f, 0x3f, 0x46)
G500       = RGBColor(0x71, 0x71, 0x7a)
G400       = RGBColor(0xa1, 0xa1, 0xaa)
G300       = RGBColor(0xd4, 0xd4, 0xd8)
G100       = RGBColor(0xf4, 0xf4, 0xf5)
LINE       = RGBColor(0xc9, 0xc9, 0xcf)
ACCENT     = RGBColor(0x43, 0x76, 0xfc)
ACCENT_INK = RGBColor(0x2f, 0x54, 0xd4)
ACCENT_BG  = RGBColor(0xe9, 0xee, 0xfe)
WHITE      = RGBColor(0xff, 0xff, 0xff)
FONT       = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = " ".join(text.split())


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def para(tf, first=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    return p


def run(p, text, size, color=INK, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.name = font
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return r


def kicker(s, text, l=0.62, t=0.5, w=12.1):
    _, tf = textbox(s, l, t, w, 0.4)
    p = para(tf, first=True)
    run(p, text.upper(), 11, G500, bold=True)


def h2(s, text, l=0.62, t=0.95, w=12.1, size=26, color=INK):
    _, tf = textbox(s, l, t, w, 0.7)
    p = para(tf, first=True)
    run(p, text, size, color, bold=True)


def card(s, l, t, w, h, title, body, accent=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT_BG if accent else WHITE
    sh.line.color.rgb = ACCENT if accent else LINE
    sh.line.width = Pt(1)
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.14)
    p = para(tf, first=True, space_after=4)
    run(p, title, 14, INK, bold=True)
    pb = para(tf, line=1.15)
    run(pb, body, 12.5, G700)
    return sh


def pill(s, l, t, text, fill=WHITE, border=LINE, color=G700):
    w = 0.16 + 0.082 * len(text)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.34))
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(0.75)
    try:
        sh.adjustments[0] = 0.5
    except Exception:
        pass
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, text, 10.5, color, bold=False)
    return w


def pill_row(s, items, t, center=True, gap=0.18, **kw):
    widths = [0.16 + 0.082 * len(x) for x in items]
    total = sum(widths) + gap * (len(items) - 1)
    x = (SW - total) / 2 if center else 0.62
    for txt, wdt in zip(items, widths):
        pill(s, x, t, txt, **kw)
        x += wdt + gap


def picture_fit(s, path, l, t, max_w, max_h, center_in=True):
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    if center_in:
        l = l + (max_w - w) / 2
    pic = s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    return pic, w, h


# ================================================================ SLIDE 1
s = slide()
picture_fit(s, os.path.join(IMG, "logo.png"), (SW - 1.0) / 2, 0.85, 1.0, 0.95, center_in=True)
_, tf = textbox(s, 1.0, 1.95, SW - 2.0, 0.4)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "TECHNISCHE VERKENNING AI-TOOLS WOO-PROCES · BZK / ECP · GOVTECHNL", 11, G500, bold=True)
_, tf = textbox(s, 1.0, 2.4, SW - 2.0, 1.1)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "Zeef", 60, INK, bold=True)
_, tf = textbox(s, 1.0, 3.65, SW - 2.0, 0.6)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "Van duizend documenten naar de honderd die ertoe doen — lokaal en navolgbaar", 18, ACCENT, bold=True)
_, tf = textbox(s, 1.0, 4.45, SW - 2.0, 0.9)
p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.5)
run(p, "Aanbieder: ", 12, INK, bold=True); run(p, "Conduction    ·    ", 12, G500)
run(p, "Licentie: ", 12, INK, bold=True); run(p, "EUPL-1.2", 12, G500)
p = para(tf, align=PP_ALIGN.CENTER); run(p, "26 juni 2026", 12, G500)
pill_row(s, ["Open source", "Volledig offline", "Boring & auditable"], 5.7)
notes(s, """Welkom. Ik ben van Conduction. Zeef doet een ding: van duizend documenten terug naar de
honderd die er werkelijk toe doen — en bij elk document kun je terugzien waarom het gekozen is.
Open source, EUPL-1.2. Wat u vandaag ziet draait volledig op een eigen machine: niets verlaat
het pand. In zeven minuten, inclusief uw vragen: hoe het werkt, lokaal versus cloud, uw
uitvraag-punten, de aanpak, en twee echte outputs — geen plaatjes.""")

# ================================================================ SLIDE 2
s = slide()
kicker(s, "Agenda")
h2(s, "Wat u in 7 minuten ziet", t=1.5, size=30)
_, tf = textbox(s, 0.62, 2.7, 9.5, 2.0)
p = para(tf, first=True, line=1.5)
run(p, "Hoe Zeef werkt · lokaal versus cloud · uitvraag-punten · de aanpak in het kort · "
        "twee echte outputs. Daarna vragenrondje.", 19, G700)
notes(s, """De rode draad: eerst begrijpen wat er ligt, dan kiezen wat telt — dat zijn twee verschillende
dingen. En: "binnen de vraag vallen" is iets anders dan "tot de belangrijkste honderd horen".
Dat onderscheid komt steeds terug. Twee slides tonen echt gegenereerde output, geen plaatjes.""")

# ================================================================ SLIDE 3
s = slide()
kicker(s, "Werkwijze · eerst begrijpen, dan kiezen")
# Big statement met accent-balk links
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.0), Inches(0.06), Inches(0.55))
bar.shadow.inherit = False; bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
_, tf = textbox(s, 0.82, 1.0, 11.0, 0.6)
p = para(tf, first=True)
run(p, "Begrijpen is niet kiezen.", 24, INK, bold=True)
picture_fit(s, os.path.join(IMG, "diamond.png"), 0.62, 2.0, 12.1, 3.3, center_in=True)
_, tf = textbox(s, 0.62, 5.6, 12.1, 1.3)
p = para(tf, first=True, line=1.45)
run(p, "Twee aparte stappen. ", 16, G700)
run(p, "Begrijpen", 16, INK, bold=True)
run(p, " brengt in kaart wat er ligt; ", 16, G700)
run(p, "kiezen", 16, INK, bold=True)
run(p, " zet de documenten op volgorde van belang voor de concrete vraag.", 16, G700)
notes(s, """Dit is de kern. Het Woo-proces vraagt twee verschillende dingen. Een: begrijpen wat er in de
berg zit — zonder zoekvraag, je kijkt breed en vat samen tot een overzicht van onderwerpen.
Twee: kiezen wat relevant is voor de concrete vraag — je begint bij alle kandidaten en houdt
de belangrijkste honderd over. De fout die je niet wilt maken: een groot onderwerp aanzien
voor het belangrijkste, of "binnen de vraag" verwarren met "in de kern". Twee stappen, nooit
door elkaar.""")

# ================================================================ SLIDE 4
s = slide()
kicker(s, "Lokaal vs cloud · en hoe zeker de keuze is")
h2(s, "Lokaal werkt — en bleek robuuster dan cloud.", t=0.95, size=24)
# Twee kaarten links-boven
card(s, 0.62, 1.7, 3.55, 1.25, "Lokaal  (vandaag getoond)",
     "Draait op een eigen machine. Niets verlaat het pand.", accent=True)
card(s, 4.35, 1.7, 3.55, 1.25, "Cloud",
     "Een krachtiger taalmodel, maar de documenten verlaten het pand.")
# Lead-tekst links
_, tf = textbox(s, 0.62, 3.15, 7.3, 3.7)
p = para(tf, first=True, line=1.4, space_after=8)
run(p, "Op grotere schaal liep het cloud-pad vast op netwerk-wachttijd, ook nadat we de last per "
        "document verlaagden. De lokale variant niet — ", 14, G700)
run(p, "lokaal bleek ook robuuster", 14, INK, bold=True)
run(p, ". We tonen vandaag de lokale selectie.", 14, G700)
p = para(tf, line=1.4, space_after=8)
run(p, "Drie manieren om te bepalen welke documenten relevant zijn — en hoe vaak ze het eens zijn:", 14, G700)
p = para(tf, line=1.4)
run(p, "De drie manieren zijn het over ", 12.5, G500)
run(p, "33 documenten", 12.5, INK, bold=True)
run(p, " eens — dat is de kern waar we het zekerst van zijn. Welke manier “beter” is, kunnen we hier "
        "niet meten: voor dit dossier bestaat geen objectieve maatstaf.", 12.5, G500)
# Venn rechts
picture_fit(s, os.path.join(IMG, "venn.png"), 8.2, 1.7, 4.9, 5.1, center_in=True)
notes(s, """Twee punten. Een, lokaal versus cloud: cloud gebruikt een krachtiger taalmodel, maar de documenten
verlaten dan het pand. En op grotere schaal liep het cloud-pad twee keer vast op netwerk-wachttijd
in de ophaal-stap, ook nadat we de last per document verlaagden. De lokale variant liep wel door —
lokaal bleek dus niet alleen veiliger maar ook robuuster. Er is vandaag dus geen cloud-selectie als
resultaat; we tonen de lokale. Twee, hoe zeker de keuze is: we lieten drie manieren los om relevantie
te bepalen — een taalmodel dat op betekenis kijkt, een klassieke methode die op woorden matcht, en een
hele eenvoudige ondergrens. Ze zijn het over drieendertig documenten eens; dat is de kern waar we het
zekerst van zijn. Welke manier beter is kunnen we hier niet hard maken — voor dit dossier is er geen
objectieve maatstaf om tegen af te meten.""")

# ================================================================ SLIDE 5
s = slide()
kicker(s, "Uw uitvraag-punten · het belangrijkste")
h2(s, "Drie punten die ertoe doen", t=1.4, size=30)
rows = [
    ("Open source", "Vrij te gebruiken en te controleren (EUPL-1.2)"),
    ("Lokaal & soeverein", "Draait in eigen huis; de documenten verlaten het pand niet"),
    ("Transparant", "Elke beslissing wordt vastgelegd en is later terug te halen"),
]
tbl_shape = s.shapes.add_table(3, 2, Inches(0.62), Inches(2.5), Inches(11.0), Inches(2.4))
tbl = tbl_shape.table
tbl.first_row = False
tbl.horz_banding = False
tbl.columns[0].width = Inches(3.3)
tbl.columns[1].width = Inches(7.7)
for i, (k, v) in enumerate(rows):
    for j, txt in enumerate((k, v)):
        cell = tbl.cell(i, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = G100 if j == 0 else WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.18)
        cell.margin_top = cell.margin_bottom = Inches(0.08)
        tf = cell.text_frame
        tf.word_wrap = True
        p = para(tf, first=True)
        run(p, txt, 16, INK, bold=(j == 0))
_, tf = textbox(s, 0.62, 5.2, 11.0, 1.2)
p = para(tf, first=True, line=1.45)
run(p, "Verder, kort: geen training op uw documenten · reeds-openbaar dossier (AVG geen beletsel) · "
        "koppelbaar aan bestaande systemen · geen kosten per verzoek. Details desgewenst.", 12.5, G500)
notes(s, """Drie kernpunten op de slide: open source onder EUPL-1.2 — vrij te gebruiken en te controleren;
lokaal en soeverein — de documenten blijven in huis; en transparant — elke beslissing wordt
vastgelegd en is terug te halen. De rest heb ik klaar als u ernaar vraagt: we trainen niet op uw
documenten; het gaat om reeds-openbaar te maken materiaal, dus AVG is geen beletsel, met een
DPIA-pad; Zeef koppelt aan bestaande systemen zoals M365, DMS en zaaksystemen; en er is geen
licentie per verzoek — open source plus eigen hardware. Schaal: vandaag is dit een verkenning,
opschalen is het pad.""")

# ================================================================ SLIDE 6
s = slide()
kicker(s, "De aanpak · in drie stappen")
h2(s, "Opruimen → kiezen → landkaart", t=0.95, size=26)
cw = 3.83
card(s, 0.62, 1.75, cw, 1.85, "1 · Opruimen",
     "Inlezen, dubbelingen samenvoegen, en documenten buiten de vraag eruit: "
     "doorstuurmails, agendaverzoeken, dubbelingen.")
card(s, 0.62 + cw + 0.3, 1.75, cw, 1.85, "2 · Kiezen",
     "Wat overblijft op volgorde van belang voor de vraag — en de ongeveer honderd "
     "belangrijkste overhouden.", accent=True)
card(s, 0.62 + 2 * (cw + 0.3), 1.75, cw, 1.85, "3 · Landkaart",
     "Een overzicht van de onderwerpen in het dossier — ook bruikbaar als er nog "
     "geen zoekvraag is.")
# Trechter
segs = [("1000", "ingelezen", False), ("1000", "leesbaar", False),
        ("888", "binnen de vraag", False), ("100", "gekozen", True)]
seg_w, gap = 1.7, 0.62
total = len(segs) * seg_w + (len(segs) - 1) * gap
x = (SW - total) / 2
ty = 4.15
for i, (val, lab, sel) in enumerate(segs):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(ty), Inches(seg_w), Inches(0.95))
    sh.shadow.inherit = False
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT_BG if sel else WHITE
    sh.line.color.rgb = ACCENT if sel else LINE; sh.line.width = Pt(1)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = para(tf, first=True, align=PP_ALIGN.CENTER, space_after=2)
    run(p, val, 22, INK, bold=True)
    p = para(tf, align=PP_ALIGN.CENTER)
    run(p, lab.upper(), 9, G500)
    if i < len(segs) - 1:
        _, atf = textbox(s, x + seg_w, ty, gap, 0.95, anchor=MSO_ANCHOR.MIDDLE)
        ap = para(atf, first=True, align=PP_ALIGN.CENTER)
        run(ap, "→", 18, G400)
    x += seg_w + gap
_, tf = textbox(s, 0.62, 5.5, 12.1, 1.2)
p = para(tf, first=True, line=1.45)
run(p, "Van 1000 documenten viel er niets af op leesbaarheid (0); 112 vielen buiten de vraag — allemaal "
        "dubbelingen — en 888 bleven over. Daarvan kozen we de 100 die het dichtst bij de vraag liggen. "
        "Waar je precies afkapt is een keuze, geen natuurwet.", 12.5, G500)
notes(s, """Drie stappen. Een, opruimen: inlezen met een vaste vingerafdruk per document zodat dubbelingen
meteen samenvallen, en documenten die buiten de vraag vallen eruit — doorstuurmails,
agendaverzoeken, dubbelingen, procesmeldingen. Dat is een spelregel, geen oordeel over belang.
Twee, kiezen: wat overblijft zetten we op volgorde van belang voor de concrete vraag en houden de
ongeveer honderd belangrijkste over. Drie, landkaart: los daarvan kan Zeef een overzicht van alle
onderwerpen maken, ook zonder zoekvraag. De trechter: duizend ingelezen, op leesbaarheid valt
niets af, honderdtwaalf vallen buiten de vraag — allemaal dubbelingen — en honderd gekozen. Waar je
precies afkapt is een keuze; het verschil rond de grens is klein, dus die honderd is een keuze, geen
natuurwet.""")

# ================================================================ SLIDE 7
s = slide()
kicker(s, "Het rapport · de echte output")
h2(s, "Wat de tool oplevert — echt, geen plaatje.", t=0.95, size=24)
# Frame rond de screenshot
_, w, h = picture_fit(s, os.path.join(IMG, "converge.png"), 0.62, 1.7, 12.1, 4.0, center_in=True)
_, tf = textbox(s, 0.62, 5.85, 12.1, 0.4)
p = para(tf, first=True)
run(p, "De top-100, met per document waarom het gekozen is — en wat eruit viel.", 12.5, G700)
run(p, "        ↗ embeds/converge-report.html", 11, ACCENT_INK, bold=True)
_, tf = textbox(s, 0.62, 6.35, 12.1, 0.4)
p = para(tf, first=True)
run(p, "Getoond: 86 (basis). Geleverd: 100 (zelfde volgorde, doorgesneden).", 11, G500)
notes(s, """Dit is geen plaatje. Dit is het rapport dat de tool zelf maakt: de gekozen honderd, bij elk
document waarom het gekozen is, hoe documenten samenhangen, en wat eruit viel en waarom. U kunt
erin navigeren. Als het op de beamer hapert, open ik het in een nieuw tabblad — de link staat
rechtsonder. Dit is de lokale run op het BZK-dossier. Een eerlijke nuance: het rapport toont er
86 als basis; de geleverde set is honderd, uit precies dezelfde volgorde, alleen op honderd
doorgesneden.""")

# ================================================================ SLIDE 8
s = slide()
kicker(s, "Buiten de vraag + de onderwerpenkaart")
h2(s, "Buiten de vraag is een spelregel, geen oordeel over belang.", t=0.95, size=21)
_, tf = textbox(s, 0.62, 1.6, 12.1, 0.5)
p = para(tf, first=True, line=1.35)
run(p, "Wat eruit gaat, gaat eruit om een ", 14, G700)
run(p, "duidelijke reden", 14, INK, bold=True)
run(p, " — en dat is altijd terug te zien:", 14, G700)
pill_row(s, ["doorstuurmails", "agendaverzoeken", "dubbelingen", "procesmeldingen",
             "eerdere mails in een thread"], 2.15, center=False, fill=G100, border=G400)
_, tf = textbox(s, 0.62, 2.75, 12.1, 0.5)
p = para(tf, first=True, line=1.35)
run(p, "Nog geen zoekvraag? ", 14, INK, bold=True)
run(p, "Dan maakt Zeef eerst een overzicht van alle onderwerpen in het dossier.", 14, G700)
picture_fit(s, os.path.join(IMG, "discover.png"), 0.62, 3.4, 12.1, 3.4, center_in=True)
_, tf = textbox(s, 0.62, 6.9, 12.1, 0.4)
p = para(tf, first=True)
run(p, "Overzicht van de onderwerpen in het dossier.", 11, G700)
run(p, "        ↗ embeds/discover-report.html", 10.5, ACCENT_INK, bold=True)
notes(s, """Buiten de vraag betekent niet onbelangrijk — het betekent: valt buiten de spelregels van dit
verzoek. De redenen zijn expliciet: doorstuurmails, agendaverzoeken, dubbelingen, procesmeldingen,
eerdere mails in een thread. Elke uitsluiting krijgt een reden mee en is terug te halen. En als de
zoekvraag nog onbekend is, maakt Zeef eerst een overzicht van alle onderwerpen in het dossier. Ook
hier een echte output, met een link rechtsonder: dat overzicht over ditzelfde BZK-dossier, met
hetzelfde lokale model als de selectie.""")

# ================================================================ SLIDE 9
s = slide()
_, tf = textbox(s, 1.0, 1.7, SW - 2.0, 1.0)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "Vragen?", 52, INK, bold=True)
_, tf = textbox(s, 2.0, 3.0, SW - 4.0, 0.6)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, "Verkenning vandaag — ", 18, G700)
run(p, "grotere schaal is het pad", 18, INK, bold=True)
run(p, ".", 18, G700)
_, tf = textbox(s, 2.0, 3.7, SW - 4.0, 1.2)
p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.4)
run(p, "Zeef aan te sluiten op bestaande systemen — ", 14, G500)
run(p, "DMS en zaaksystemen", 14, INK, bold=True)
run(p, " — en in de lokale variant op ", 14, G500)
run(p, "Nextcloud", 14, INK, bold=True)
run(p, " en ", 14, G500)
run(p, "OpenAnonymiser", 14, INK, bold=True)
run(p, " voor de AVG-voorbewerking (contextverrijkend anonimiseren i.p.v. ", 14, G500)
run(p, "lakken", 14, INK, bold=True)
run(p, ").", 14, G500)
_, tf = textbox(s, 2.0, 5.05, SW - 4.0, 0.8)
p = para(tf, first=True, align=PP_ALIGN.CENTER, line=1.5)
run(p, "Zeef", 12, INK, bold=True); run(p, " · Conduction  ·  EUPL-1.2", 12, G500)
p = para(tf, align=PP_ALIGN.CENTER); run(p, "mark@conduction.nl", 12, G500)
pill_row(s, ["Open source", "Soeverein & offline", "Elke beslissing herleidbaar"], 6.1)
notes(s, """Samengevat: Zeef brengt u van duizend documenten terug naar de honderd die ertoe doen, en bij elke
stap kunt u terugzien waarom. Lokaal waar het moet, elke beslissing terug te halen. Open source,
EUPL-1.2. Het sluit aan op uw bestaande systemen — DMS en zaaksystemen — en in de lokale variant op
Nextcloud en OpenAnonymiser voor de AVG-voorbewerking. Wat u vandaag zag is een verkenning;
opschalen is het pad. Dank — ik beantwoord graag uw vragen. Contact: mark apenstaartje conduction
punt nl.""")

prs.save(OUT)
print("Geschreven:", os.path.normpath(OUT), "·", len(prs.slides._sldIdLst), "slides")
